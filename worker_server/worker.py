import sys
import sqlite3
import os
import time
import io
import string
import itertools
import zipfile
from concurrent.futures import ThreadPoolExecutor

try:
    import fitz # PyMuPDF
    from PIL import Image
except ImportError:
    pass

# --- Dependencies Check ---
try:
    import pikepdf
except ImportError:
    print("Error: pikepdf not installed")
    sys.exit(1)

# Argument parsing
if len(sys.argv) < 5:
    print("Usage: python worker.py <input> <output> <job_id> <task_type>")
    sys.exit(1)

INPUT_PATH = sys.argv[1]
OUTPUT_PATH = sys.argv[2]
JOB_ID = sys.argv[3]
TASK_TYPE = sys.argv[4]
DB_PATH = os.path.join(os.path.dirname(__file__), 'db', 'tasks.sqlite')

# --- DB Helper ---
def update_db(status, password=None):
    try:
        conn = sqlite3.connect(DB_PATH, timeout=10)
        c = conn.cursor()
        c.execute("UPDATE tasks SET status=?, password_found=?, updated_at=(DATETIME('now', 'localtime')) WHERE id=?", 
                  (status, password, JOB_ID))
        conn.commit()
        conn.close()
    except Exception as e:
        print(f"DB Error: {e}")

def check_abandonment():
    """Detect if the user closed the web page (Stopped polling)"""
    try:
        conn = sqlite3.connect(DB_PATH, timeout=5)
        c = conn.cursor()
        # threshold: 600 seconds of inactivity or explicit abandoned status
        c.execute("PRAGMA journal_mode=WAL;")
        c.execute("""
            SELECT CASE 
                WHEN status = 'abandoned' THEN 1
                WHEN (julianday('now') - julianday(last_heartbeat)) * 86400 > 600 THEN 1 
                ELSE 0 
            END FROM tasks WHERE id=?
        """, (JOB_ID,))
        res = c.fetchone()
        conn.close()
        return bool(res[0]) if res else False
    except:
        return False

def check_queue():
    """Wait if another unlock task is already processing (Sequential Queue for Unlock)"""
    if TASK_TYPE.lower() != 'unlock':
        return
    
    print(f"⏳ Task {JOB_ID} checking queue for 'unlock' missions...")
    # Stay in pending while waiting
    
    while True:
        try:
            conn = sqlite3.connect(DB_PATH, timeout=15)
            c = conn.cursor()
            # 1. Check if I am the oldest ACTIVE pending or processing unlock task
            # We ignore tasks that haven't had a heartbeat for 600+ seconds
            c.execute("PRAGMA journal_mode=WAL;")
            c.execute("""
                SELECT id FROM tasks 
                WHERE type='unlock' 
                AND status IN ('processing', 'pending')
                AND (julianday('now') - julianday(last_heartbeat)) * 86400 < 600
                ORDER BY created_at ASC LIMIT 1
            """)
            res = c.fetchone()
            conn.close()
            
            if res and res[0] == JOB_ID:
                print(f"✅ It is my turn ({JOB_ID}). Starting unlock.")
                update_db('processing')
                break
            elif not res:
                # No active tasks found? (Shouldn't happen to me, but safe check)
                update_db('processing')
                break
            else:
                # Someone else is active and older than me
                if check_abandonment():
                    print(f"🛑 Task {JOB_ID} cancelled while in queue.")
                    update_db('abandoned')
                    os._exit(0)
                time.sleep(10)
        except Exception as e:
            print(f"Queue check error: {e}")
            time.sleep(5)

# ==========================================
# 🔐 MODE: UNLOCK (Hyper-Speed Neural v19)
# ==========================================
def run_unlock():
    print(f"🚀 Initializing Hyper-Speed Node: {JOB_ID}")
    
    try:
        with open(INPUT_PATH, 'rb') as f:
            FILE_BYTES = f.read()
    except Exception as e:
        print(f"Read Error: {e}")
        update_db('failed')
        return

    mem_buf = io.BytesIO(FILE_BYTES)
    
    # 0. Pre-check: Is it actually encrypted?
    try:
        with pikepdf.open(mem_buf) as pdf:
            # Opened without password! 
            print("✨ File is NOT encrypted. Saving directly...")
            pdf.save(OUTPUT_PATH)
            update_db('completed', 'No Password Required')
            return
    except pikepdf.PasswordError:
        print("🔒 File is encrypted. Starting Decryption...")
    except Exception as e:
        print(f"File Error: {e}")
        update_db('failed')
        return

    found_password = None
    stop_event = False

    def check_solution(pwd):
        if not pwd: return False
        try:
            # 1. Primary Check (Pikepdf)
            mem_buf.seek(0)
            with pikepdf.open(mem_buf, password=pwd) as pdf:
                if len(pdf.pages) == 0: return False
                
                # 2. Secondary Strict Check (PyMuPDF) - CRITICAL for preventing False Positives
                # Only run if fitz is available (it should be)
                if 'fitz' in sys.modules:
                    try:
                        mem_buf.seek(0)
                        # Re-open stream with fitz for authentication check
                        with fitz.open(stream=mem_buf.read(), filetype="pdf") as doc:
                            if not doc.authenticate(pwd):
                                # Pikepdf said YES, but Fitz said NO -> Likely False Positive (Collision)
                                return False
                    except Exception as e:
                        # If fitz fails to open, it's definitely bad
                        return False
                
                pdf.save(OUTPUT_PATH)
                return True
        except:
            return False

    def solver_wrapper(candidate_gen, name="Node"):
        nonlocal found_password, stop_event
        print(f"  [+] {name} spinning up...")
        for pwd in candidate_gen:
            if stop_event: return None
            if check_solution(pwd):
                stop_event = True
                return pwd
        return None

    # --- TOP PRIORITY 1: Years and Pure Numbers (Fastest) ---
    def gen_priority_years_and_nums():
        # 1. Buddhist Era (พ.ศ.) & CE Years (Immediate Strike)
        curr = 2026
        for y in range(curr + 5, 1950, -1):
            yield str(y + 543) # BE first (พ.ศ.)
            yield str(y)       # CE
            
        # 2. Common Digital Pins (4-digit and 6-digit)
        for i in range(1000000):
            if i < 10000: yield f"{i:04d}"
            yield f"{i:06d}"

    # --- NEW TOP PRIORITY: Smart Date-of-Birth (All Formats) ---
    def gen_priority_dates():
        import calendar
        # Focus years: 1940-2027 (CE) and 2483-2570 (BE)
        for year_ce in range(2027, 1940, -1):
            year_be = year_ce + 543
            for y in [year_ce, year_be]:
                y_str = str(y)
                y_short = y_str[2:]
                for month in range(1, 13):
                    _, last_day = calendar.monthrange(year_ce, month)
                    for day in range(1, last_day + 1):
                        d, m = f"{day:02d}", f"{month:02d}"
                        # 1. Standard Pure Numeric (Most Common)
                        yield f"{d}{m}{y_str}"   # DDMMYYYY
                        yield f"{y_str}{m}{d}"   # YYYYMMDD
                        yield f"{d}{m}{y_short}" # DDMMYY
                        
                        # 2. With Separators (As requested: dd/mm/yyyy)
                        yield f"{d}/{m}/{y_str}"
                        yield f"{m}/{d}/{y_str}"
                        yield f"{y_str}/{m}/{d}"
                        yield f"{d}-{m}-{y_str}"
                        yield f"{y_str}-{m}-{d}"

    def gen_priority_nums():
        # Quick Pins
        for i in range(10000): yield f"{i:04d}"
        for i in range(1000000): yield f"{i:06d}"
        # Mobile Numbers (Typical TH 10 digits)
        # We can add common prefixes if needed

    def gen_thai_hybrid():
        names = ["moon", "boss", "manky", "admin", "password", "12345678"]
        years = [str(y) for y in range(2027, 1980, -1)] + [str(y+543) for y in range(2027, 1980, -1)]
        for n in names:
            for y in years:
                yield f"{n}{y}"
                yield f"{y}{n}"
                yield f"{n}@{y}"

    # --- FALLBACK: Heavy Deep Scan (Inform user this may take hours) ---
    def gen_deep_numeric():
        for i in range(1000000, 100000000): # 7-8 digits
            yield str(i)

    # --- FALLBACK: Exhaustive 8-Digit and Long Numeric ---
    def gen_heavy_numeric():
        for i in range(1000000, 100000000): # 7-8 digits
            yield str(i)

    # Execution Deployment
    # Optimize: Use CPU Count minus 1 to keep the web server responsive
    MAX_WORKERS = max(1, (os.cpu_count() or 2) - 1)
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as exe:
        futures = []
        # STEP 1: Priority Scan (DOB, Common Pins) - FAST
        futures.append(exe.submit(solver_wrapper, gen_priority_dates(), "BirthDateStrike"))
        futures.append(exe.submit(solver_wrapper, gen_priority_nums(), "CommonPins"))
        futures.append(exe.submit(solver_wrapper, gen_thai_hybrid(), "ThaiHybrid"))
        
        heartbeat_counter = 0
        priority_done = False

        while any(not f.done() for f in futures):
            if stop_event: break
            
            # Check for abandonment every ~5 seconds
            heartbeat_counter += 1
            if heartbeat_counter >= 12:
                heartbeat_counter = 0
                if check_abandonment():
                    print(f"🛑 Task {JOB_ID} ABANDONED. Forcing immediate exit.")
                    update_db('abandoned')
                    os._exit(0)
            
            # If all priority scans are done and no password found, signal Deep Scan
            if not priority_done and all(f.done() for f in futures):
                priority_done = True
                found_yet = any(f.result() for f in futures if f.done())
                if not found_yet:
                    print("⚠️ Priority scans failed. Entering DEEP SCAN PHASE...")
                    # Notify DB that it's going to take a long time
                    update_db('processing') # You might want a custom detail or just let it run
                    # Add NEW Heavy Futures
                    futures.append(exe.submit(solver_wrapper, gen_deep_numeric(), "DeepScan-Numeric"))

            time.sleep(0.4)

        for f in futures:
            if f.done():
                try:
                    res = f.result()
                    if res: found_password = res; break
                except: pass

    if found_password:
        update_db('completed', found_password)
        print(f"✅ FOUND: {found_password}")
    else:
        update_db('failed')
        print("❌ Search exhausted.")

# ==========================================
# 📊 MODE: PDF TO PPT (Image Based)
# ==========================================
def run_pdf_to_ppt():
    print(f"📊 Starting PDF->PPT Task: {JOB_ID}")
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Missing libraries: pdf2image or python-pptx")
        update_db('failed')
        return

    try:
        # Convert PDF to Images
        # Note: 'poppler' must be installed on the system (apt install poppler-utils)
        images = convert_from_path(INPUT_PATH)
        
        prs = Presentation()
        # Clean default slides
        
        for i, image in enumerate(images):
            # Check for abandonment every few pages
            if i % 5 == 0 and check_abandonment():
                print(f"🛑 PDF->PPT Task {JOB_ID} ABANDONED.")
                update_db('abandoned')
                os._exit(0)

            # Save image temporarily
            img_path = f"{OUTPUT_PATH}_temp_{i}.jpg"
            image.save(img_path, 'JPEG')
            
            # Create blank slide
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to slide (Fit to page logic could be added)
            # Default: specific size or fit
            left = top = Inches(0)
            # A4 aspect ratio approx
            slide.shapes.add_picture(img_path, left, top, height=Inches(7.5)) 
            
            os.remove(img_path) # Cleanup temp image

        prs.save(OUTPUT_PATH)
        update_db('completed')
        
    except Exception as e:
        print(f"PPT Conversion Error: {e}")
        update_db('failed')

# ==========================================
# 📈 MODE: PDF TO EXCEL
# ==========================================
def run_pdf_to_excel():
    print(f"📈 Starting PDF->Excel Task: {JOB_ID}")
    try:
        import pdfplumber
        import pandas as pd
    except ImportError:
        print("Missing libraries: pdfplumber or pandas")
        update_db('failed')
        return

    try:
        all_tables = []
        with pdfplumber.open(INPUT_PATH) as pdf:
            for i, page in enumerate(pdf.pages):
                # Check for abandonment
                if i % 10 == 0 and check_abandonment():
                    print(f"🛑 PDF->Excel Task {JOB_ID} ABANDONED.")
                    update_db('abandoned')
                    os._exit(0)

                tables = page.extract_tables()
                for table in tables:
                    df = pd.DataFrame(table)
                    all_tables.append(df)
        
        if all_tables:
            # Concatenate all tables or create multiple sheets?
            # Simple approach: Concat all
            final_df = pd.concat(all_tables)
            final_df.to_excel(OUTPUT_PATH, index=False, header=False)
            update_db('completed')
        else:
            print("No tables found")
            # Create empty excel
            pd.DataFrame(["No tables found"]).to_excel(OUTPUT_PATH)
            update_db('completed') # Treated as success but empty
            
    except Exception as e:
        print(f"Excel Conversion Error: {e}")
        update_db('failed')


# ==========================================
# 📄 MODE: PPT TO PDF (LibreOffice)
# ==========================================
def run_ppt_to_pdf():
    print(f"📄 Starting PPT->PDF Task: {JOB_ID}")
    import subprocess
    
    try:
        # Check if LibreOffice is installed
        # Common names: libreoffice, soffice
        cmd = ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(OUTPUT_PATH), INPUT_PATH]
        
        print(f"Executing: {' '.join(cmd)}")
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        if result.returncode == 0:
            # LibreOffice output name logic: same name as input but .pdf
            # We need to ensure it matches OUTPUT_PATH expected by system
            
            # Theoretical output
            base_name = os.path.splitext(os.path.basename(INPUT_PATH))[0]
            generated_pdf = os.path.join(os.path.dirname(OUTPUT_PATH), base_name + '.pdf')
            
            if os.path.exists(generated_pdf):
                # Rename to expected OUTPUT_PATH if needed
                if generated_pdf != OUTPUT_PATH:
                     if os.path.exists(OUTPUT_PATH): os.remove(OUTPUT_PATH)
                     os.rename(generated_pdf, OUTPUT_PATH)
                
                update_db('completed')
            else:
                print("LibreOffice finished but output file missing.")
                print(f"Expected: {generated_pdf}")
                update_db('failed')
        else:
            print(f"LibreOffice Error: {result.stderr.decode()}")
            update_db('failed')

    except Exception as e:
        print(f"PPT->PDF Error: {e}")
# ==========================================
# 📝 MODE: PDF TO WORD
# ==========================================
def run_pdf_to_word():
    print(f"📝 Starting PDF->Word Task: {JOB_ID}")
    try:
        from pdf2docx import Converter
    except ImportError:
        print("Missing library: pdf2docx")
        update_db('failed')
        return

    try:
        cv = Converter(INPUT_PATH)
        cv.convert(OUTPUT_PATH, start=0, end=None)
        cv.close()
        update_db('completed')
    except Exception as e:
        print(f"PDF->Word Error: {e}")
        update_db('failed')

# ==========================================
# 📄 MODE: WORD TO PDF (LibreOffice)
# ==========================================
def run_word_to_pdf():
    print(f"📄 Starting Word->PDF Task: {JOB_ID}")
    import subprocess
    
    try:
        # LibreOffice Handles .doc and .docx
        cmd = ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(OUTPUT_PATH), INPUT_PATH]
        
        print(f"Executing: {' '.join(cmd)}")
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        if result.returncode == 0:
            # Logic same as PPT->PDF
            base_name = os.path.splitext(os.path.basename(INPUT_PATH))[0]
            generated_pdf = os.path.join(os.path.dirname(OUTPUT_PATH), base_name + '.pdf')
            
            if os.path.exists(generated_pdf):
                if generated_pdf != OUTPUT_PATH:
                     if os.path.exists(OUTPUT_PATH): os.remove(OUTPUT_PATH)
                     os.rename(generated_pdf, OUTPUT_PATH)
                update_db('completed')
            else:
                update_db('failed')
        else:
            print(f"LibreOffice Error: {result.stderr.decode()}")
            update_db('failed')

    except Exception as e:
        print(f"Word->PDF Error: {e}")
        update_db('failed')

# ==========================================
# 🖼️ MODE: PDF TO IMAGE (ZIP)
# ==========================================
def run_pdf_to_image():
    print(f"🖼️ Starting PDF->Image Task: {JOB_ID}")
    try:
        doc = fitz.open(INPUT_PATH)
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w") as zip_file:
            for page_index in range(len(doc)):
                # Check for abandonment
                if page_index % 10 == 0 and check_abandonment():
                    print(f"🛑 PDF->Img Task {JOB_ID} ABANDONED.")
                    update_db('abandoned')
                    os._exit(0)

                page = doc.load_page(page_index)
                pix = page.get_pixmap()
                img_data = pix.tobytes("jpg")
                zip_file.writestr(f"page_{page_index+1}.jpg", img_data)
        
        with open(OUTPUT_PATH, "wb") as f:
            f.write(zip_buffer.getvalue())
        update_db('completed')
    except Exception as e:
        print(f"PDF->Image Error: {e}")
        update_db('failed')

# ==========================================
# 📄 MODE: IMAGE TO PDF
# ==========================================
def run_image_to_pdf():
    print(f"📄 Starting Image->PDF Task: {JOB_ID}")
    try:
        # Check if input is ZIP (Multiple images)
        if zipfile.is_zipfile(INPUT_PATH):
            print("📦 Multiple images detected (ZIP). Merging...")
            images_list = []
            with zipfile.ZipFile(INPUT_PATH, 'r') as zip_ref:
                # Extract to temp directory
                temp_dir = os.path.join(os.path.dirname(INPUT_PATH), f"tmp_{JOB_ID}")
                os.makedirs(temp_dir, exist_ok=True)
                zip_ref.extractall(temp_dir)
                
                # Get all files and sort them (Django sent them with image_001_ prefix)
                extracted_files = sorted([os.path.join(temp_dir, f) for f in os.listdir(temp_dir) if os.path.isfile(os.path.join(temp_dir, f))])
                
                for img_path in extracted_files:
                    try:
                        img = Image.open(img_path)
                        if img.mode != 'RGB':
                            img = img.convert('RGB')
                        images_list.append(img)
                    except:
                        continue # Skip non-image files
                
                if images_list:
                    # Save all images into one PDF
                    images_list[0].save(OUTPUT_PATH, "PDF", resolution=100.0, save_all=True, append_images=images_list[1:])
                    print(f"✅ Created PDF with {len(images_list)} pages.")
                else:
                    raise Exception("No valid images found in ZIP")
                
                # Cleanup temp
                import shutil
                shutil.rmtree(temp_dir)
        else:
            # Single Image
            print("🖼️ Single image detected.")
            img = Image.open(INPUT_PATH)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            img.save(OUTPUT_PATH, "PDF", resolution=100.0)
            
        update_db('completed')
    except Exception as e:
        print(f"Image->PDF Error: {e}")
        update_db('failed')

# ==========================================
# 📏 MODE: IMAGE RESIZE
# ==========================================
def run_image_resize():
    print(f"📏 Starting Image Resize Task: {JOB_ID}")
    try:
        img = Image.open(INPUT_PATH)
        # Resize to 50% or max 1280px
        width, height = img.size
        new_width = min(width // 2, 1280)
        ratio = new_width / width
        new_height = int(height * ratio)
        
        img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
        img.save(OUTPUT_PATH, quality=85, optimize=True)
        update_db('completed')
    except Exception as e:
        print(f"Image Resize Error: {e}")
        update_db('failed')

# ==========================================
# 🔄 MODE: IMAGE CONVERT
# ==========================================
def run_image_convert():
    print(f"🔄 Starting Image Convert Task: {JOB_ID}")
    try:
        img = Image.open(INPUT_PATH)
        # Default convert to JPG for broad compatibility, or handle via name
        if OUTPUT_PATH.lower().endswith('.webp'):
            img.save(OUTPUT_PATH, "WEBP", quality=80)
        elif OUTPUT_PATH.lower().endswith('.png'):
            img.save(OUTPUT_PATH, "PNG")
        else:
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            img.save(OUTPUT_PATH, "JPEG", quality=90)
        update_db('completed')
    except Exception as e:
        print(f"Image Convert Error: {e}")
        update_db('failed')

# ==========================================
# 🎮 MAIN DISPATCHER
# ==========================================
if __name__ == '__main__':
    # Mark as pending initially to let check_queue decide when to 'process'
    update_db('pending')
    
    # Sequential Queue Check (Specifically for Unlock)
    # This will wait until it is our turn, then update status to 'processing'
    check_queue()
    
    # If not unlock, it will skip queue and we should set processing now
    if TASK_TYPE.lower() != 'unlock':
        update_db('processing')
    
    mapping = {
        'unlock': run_unlock,
        'pdf-to-ppt': run_pdf_to_ppt,
        'pdf-to-excel': run_pdf_to_excel,
        'ppt-to-pdf': run_ppt_to_pdf,
        'pdf-to-word': run_pdf_to_word,
        'word-to-pdf': run_word_to_pdf,
        'pdf-to-image': run_pdf_to_image,
        'image-to-pdf': run_image_to_pdf,
        'image-resize': run_image_resize,
        'image-convert': run_image_convert
    }

    if TASK_TYPE in mapping:
        mapping[TASK_TYPE]()
    else:
        print(f"Unknown Task Type: {TASK_TYPE}")
        update_db('failed')
